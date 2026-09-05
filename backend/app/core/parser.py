"""多格式解析器：PDF / DOCX / PPTX / EPUB / TXT / MD → [(locator, text), ...]。

设计要点：
- 引用粒度跟格式走：PDF=页码、DOCX=段落序号、PPTX=幻灯片号、
  EPUB=章节名、TXT/MD=行号段；
- 扫描版 PDF（提取不到文本）抛 ScannedPdfError，由上层标记 rejected
  并向用户明示拒收原因，而不是静默产出空索引。
"""

from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader


class ScannedPdfError(Exception):
    """扫描版 PDF：无法提取文本，需明确拒收。"""


@dataclass
class ParsedUnit:
    """解析产物单元：定位符 + 文本。"""

    locator: str  # 如 "第12页" / "第18张幻灯片" / "第3章 图与图的基本概念" / "段落45"
    text: str


# 各格式对应的定位符类型（写入 documents.locator_type）
LOCATOR_TYPES = {
    "pdf": "page",
    "docx": "paragraph",
    "pptx": "slide",
    "epub": "section",
    "txt": "line",
    "md": "line",
}


def parse_pdf(path: Path) -> list[ParsedUnit]:
    """PDF：逐页提取文本；检测扫描件（全部页文本量近零）。"""
    reader = PdfReader(str(path))
    units: list[ParsedUnit] = []
    total_chars = 0
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        total_chars += len(text)
        if text:
            units.append(ParsedUnit(locator=f"第{i}页", text=text))
    # 扫描件判定：每页平均字符数过低说明文本层基本不存在
    if reader.pages and total_chars < len(reader.pages) * 10:
        raise ScannedPdfError(
            "该 PDF 几乎不含可选文本（疑似扫描版/图片版），暂不支持，请提供文字版文件"
        )
    return units


def parse_docx(path: Path) -> list[ParsedUnit]:
    """DOCX：按段落提取；空段落跳过，定位符为段落序号。"""
    from docx import Document

    doc = Document(str(path))
    units = []
    for i, para in enumerate(doc.paragraphs, start=1):
        text = para.text.strip()
        if text:
            units.append(ParsedUnit(locator=f"段落{i}", text=text))
    return units


def parse_pptx(path: Path) -> list[ParsedUnit]:
    """PPTX：每张幻灯片的全部文本框合并为一个单元。"""
    from pptx import Presentation

    prs = Presentation(str(path))
    units = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
        if lines:
            units.append(ParsedUnit(locator=f"第{i}张幻灯片", text="\n".join(lines)))
    return units


def parse_epub(path: Path) -> list[ParsedUnit]:
    """EPUB：按章节（spine 顺序的 HTML 文档）提取纯文本。

    每章先尝试从标题标签提取章节名，取不到则用"第 N 节"。
    常见中文乱码问题：BeautifulSoup 解析后转纯文本处理。
    """
    import bs4
    from ebooklib import ITEM_DOCUMENT, epub

    book = epub.read_epub(str(path))
    units = []
    idx = 0
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = bs4.BeautifulSoup(item.get_content(), "html.parser")
        # 章节名：优先 h1/h2，其次 title 标签
        title_tag = soup.find(["h1", "h2"]) or soup.find("title")
        title = title_tag.get_text(strip=True) if title_tag else ""
        text = soup.get_text("\n", strip=True)
        # 去掉与标题重复的首行，避免分块后重复内容
        if title and text.startswith(title):
            text = text[len(title) :].strip()
        if text:
            idx += 1
            locator = title if title else f"第{idx}节"
            units.append(ParsedUnit(locator=locator[:60], text=text))
    return units


def parse_txt(path: Path) -> list[ParsedUnit]:
    """TXT/MD：按空行分段，定位符为起始行号（MD 的 # 标题作为切分增强）。"""
    raw = path.read_text(encoding="utf-8", errors="ignore")
    lines = raw.splitlines()
    units = []
    buf: list[str] = []
    start_line = 1
    for no, line in enumerate(lines, start=1):
        is_md_heading = path.suffix.lower() == ".md" and line.lstrip().startswith("#")
        # 切分条件：空行 或 Markdown 标题行
        if (line.strip() == "" or is_md_heading) and buf:
            units.append(
                ParsedUnit(
                    locator=f"第{start_line}行",
                    text="\n".join(buf).strip(),
                )
            )
            buf = []
        if line.strip() != "":
            if not buf:
                start_line = no
            buf.append(line.rstrip())
    if buf:
        units.append(ParsedUnit(locator=f"第{start_line}行", text="\n".join(buf).strip()))
    return units


PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "epub": parse_epub,
    "txt": parse_txt,
    "md": parse_txt,  # TXT 与 MD 共用行号解析，MD 额外按标题切分
}


def parse_file(path: Path, file_type: str) -> list[ParsedUnit]:
    """按格式路由到对应解析器。未知格式直接抛错（上传层已校验扩展名）。"""
    parser = PARSERS.get(file_type)
    if parser is None:
        raise ValueError(f"不支持的文件格式：{file_type}")
    return parser(path)


def detect_scanned(path: Path) -> bool:
    """预检 PDF 是否为扫描件（用于上传早期反馈，避免白等一轮解析）。"""
    try:
        parse_pdf(path)
        return False
    except ScannedPdfError:
        return True
